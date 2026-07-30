from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = Path("docs") / "Задание_3_Мониторинг_CafeHelp_Презентация.pptx"

BG = RGBColor(248, 250, 252)
NAVY = RGBColor(22, 51, 91)
ACCENT = RGBColor(38, 122, 184)
GREEN = RGBColor(34, 197, 94)
YELLOW = RGBColor(245, 158, 11)
RED = RGBColor(239, 68, 68)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(100, 116, 139)
WHITE = RGBColor(255, 255, 255)
PANEL = RGBColor(255, 255, 255)
PANEL_BORDER = RGBColor(203, 213, 225)


def set_background(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_top_bar(slide, title, subtitle=None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.7)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.color.rgb = NAVY

    tx = slide.shapes.add_textbox(Inches(0.45), Inches(0.12), Inches(8.8), Inches(0.35))
    p = tx.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Segoe UI"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = WHITE

    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.47), Inches(0.78), Inches(12.0), Inches(0.35))
        p2 = sb.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Segoe UI"
        r2.font.size = Pt(11)
        r2.font.color.rgb = MUTED


def add_panel(slide, x, y, w, h, title=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = PANEL_BORDER
    shape.line.width = Pt(1)
    if title:
        tx = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.08), w - Inches(0.36), Inches(0.25))
        p = tx.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name = "Segoe UI"
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = NAVY
    return shape


def add_bullets(slide, x, y, w, h, items, font_size=20, color=TEXT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.font.name = "Segoe UI"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return tb


def add_text(slide, x, y, w, h, text, font_size=20, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Segoe UI"
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_metric_card(slide, x, y, w, h, title, text, accent_color):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = accent_color
    card.line.width = Pt(2)

    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Inches(0.16), h)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent_color
    stripe.line.color.rgb = accent_color

    add_text(slide, x + Inches(0.28), y + Inches(0.15), w - Inches(0.4), Inches(0.3), title, 14, True, NAVY)
    add_text(slide, x + Inches(0.28), y + Inches(0.48), w - Inches(0.4), h - Inches(0.55), text, 11, False, TEXT)


def add_table_like(slide, x, y, col_widths, headers, rows):
    row_h = Inches(0.52)
    cur_x = x
    for idx, header in enumerate(headers):
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, cur_x, y, col_widths[idx], row_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(219, 234, 254)
        rect.line.color.rgb = PANEL_BORDER
        add_text(slide, cur_x + Inches(0.05), y + Inches(0.06), col_widths[idx] - Inches(0.1), row_h - Inches(0.1), header, 10, True, NAVY)
        cur_x += col_widths[idx]

    for row_idx, row in enumerate(rows):
        cur_x = x
        top = y + row_h * (row_idx + 1)
        fill = WHITE if row_idx % 2 == 0 else RGBColor(248, 250, 252)
        for idx, value in enumerate(row):
            rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, cur_x, top, col_widths[idx], row_h)
            rect.fill.solid()
            rect.fill.fore_color.rgb = fill
            rect.line.color.rgb = PANEL_BORDER
            add_text(slide, cur_x + Inches(0.05), top + Inches(0.05), col_widths[idx] - Inches(0.1), row_h - Inches(0.08), value, 9, False, TEXT)
            cur_x += col_widths[idx]


def slide_title(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RGBColor(241, 245, 249))

    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.6), Inches(12.2), Inches(5.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = PANEL_BORDER

    add_text(slide, Inches(0.95), Inches(1.15), Inches(10.8), Inches(0.8), title, 28, True, NAVY)
    add_text(slide, Inches(0.98), Inches(2.05), Inches(10.4), Inches(1.6), subtitle, 18, False, TEXT)
    add_text(slide, Inches(0.98), Inches(5.55), Inches(4.5), Inches(0.3), "Проект: CafeHelp", 12, True, ACCENT)
    add_text(slide, Inches(0.98), Inches(5.9), Inches(7.0), Inches(0.3), "Система автоматизации кафе: касса, смены, склад, печать и аналитика", 11, False, MUTED)


def slide_goal(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_bar(slide, "Цель и контекст мониторинга", "Почему для CafeHelp нужны метрики и система контроля")

    add_panel(slide, Inches(0.45), Inches(1.25), Inches(6.15), Inches(5.7), "Критичные функции сервиса")
    add_bullets(
        slide,
        Inches(0.7), Inches(1.75), Inches(5.6), Inches(4.8),
        [
            "приём и создание заказов в кассе",
            "печать чеков и заказов на кухню",
            "открытие и закрытие смен",
            "складские остатки и движения",
            "операционная аналитика по заказам",
        ],
        font_size=20,
    )

    add_panel(slide, Inches(6.8), Inches(1.25), Inches(6.05), Inches(5.7), "Цель мониторинга")
    add_bullets(
        slide,
        Inches(7.05), Inches(1.75), Inches(5.45), Inches(4.8),
        [
            "обнаруживать сбои до того, как они влияют на выручку",
            "контролировать стабильность и скорость критичных сценариев",
            "снижать риск потери заказов и простоев смены",
            "поддерживать SLA и ускорять восстановление после инцидентов",
        ],
        font_size=20,
    )


def slide_metrics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_bar(slide, "Ключевые метрики сервиса", "Набор метрик, отражающих качество и стабильность CafeHelp")

    cards = [
        ("Доступность", "Доля успешных проверок критичных функций order, shift, warehouse, print.", ACCENT),
        ("P95 отклика", "95-й перцентиль времени ответа для ключевых API сценариев.", ACCENT),
        ("Уровень ошибок", "Доля запросов с 5xx или неуспешным бизнес-результатом.", RED),
        ("Успешность печати", "Показывает, дошёл ли заказ до принтера и кухни.", GREEN),
        ("Доля задержек", "Процент заказов, где была зафиксирована задержка исполнения.", YELLOW),
        ("Критичные остатки", "Число ингредиентов, остаток которых ниже минимального порога.", RED),
        ("MTTR", "Среднее время восстановления после инцидента.", YELLOW),
    ]

    positions = [
        (Inches(0.45), Inches(1.25)),
        (Inches(4.45), Inches(1.25)),
        (Inches(8.45), Inches(1.25)),
        (Inches(0.45), Inches(3.2)),
        (Inches(4.45), Inches(3.2)),
        (Inches(8.45), Inches(3.2)),
        (Inches(0.45), Inches(5.15)),
    ]

    for (title, text, color), (x, y) in zip(cards, positions):
        add_metric_card(slide, x, y, Inches(3.55), Inches(1.55), title, text, color)

    add_panel(slide, Inches(4.45), Inches(5.05), Inches(7.55), Inches(1.75), "Почему эти метрики важны")
    add_text(
        slide,
        Inches(4.7), Inches(5.5), Inches(7.0), Inches(1.0),
        "Они связывают техническое состояние сервиса с операционными рисками кафе: остановкой продаж, потерей заказов, задержками кухни и дефицитом ингредиентов.",
        15, False, TEXT
    )


def slide_method(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_bar(slide, "Методика измерения", "Источники данных, расчёты и периодичность контроля")

    add_panel(slide, Inches(0.4), Inches(1.2), Inches(12.45), Inches(5.95), None)
    headers = ["Метрика", "Источник", "Расчёт", "Периодичность", "Ограничение"]
    widths = [Inches(2.0), Inches(2.55), Inches(3.25), Inches(1.55), Inches(3.1)]
    rows = [
        ["Доступность", "health-check, логи", "успешные / все * 100%", "1 мин", "не ловит UX-деградацию"],
        ["P95 отклика", "access-логи, APM", "p95 по критичным API", "5 мин", "зависит от нагрузки"],
        ["Уровень ошибок", "HTTP и exception логи", "ошибки / все * 100%", "5 мин", "важно отделять user/system"],
        ["Успешность печати", "логи print jobs", "успешные / все * 100%", "5 мин", "зависит от драйвера и ОС"],
        ["Доля задержек", "таблица заказов", "delay > 0 / все * 100%", "смена/день", "влияет и кухня"],
        ["Критичные остатки", "остатки + движения", "count(qty <= threshold)", "15 мин", "порог зависит от unit"],
        ["MTTR", "журнал инцидентов", "avg(resolve - open)", "неделя", "нужна фиксация инцидентов"],
    ]
    add_table_like(slide, Inches(0.6), Inches(1.5), widths, headers, rows)


def slide_dashboard(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_bar(slide, "Макет дашборда", "Структура экрана мониторинга для операторов, администраторов и менеджмента")

    add_panel(slide, Inches(0.6), Inches(1.35), Inches(12.1), Inches(5.4), "Схема дашборда")

    add_metric_card(slide, Inches(1.0), Inches(1.95), Inches(10.8), Inches(0.75), "Верхняя зона KPI", "SLA / Availability | P95 API | Error Rate | Print Success | MTTR", ACCENT)
    add_metric_card(slide, Inches(1.0), Inches(2.95), Inches(10.8), Inches(0.75), "Операционная зона", "Delayed Orders | Critical Stock | Active Alerts | Incidents Today", YELLOW)
    add_metric_card(slide, Inches(1.0), Inches(3.95), Inches(5.2), Inches(0.95), "Графики", "Latency/Error Trend (24h)", GREEN)
    add_metric_card(slide, Inches(6.6), Inches(3.95), Inches(5.2), Inches(0.95), "Графики", "Delayed Orders / Print Success Trend", GREEN)
    add_metric_card(slide, Inches(1.0), Inches(5.15), Inches(10.8), Inches(0.75), "Лента инцидентов", "время | severity | метрика | объект | ответственный | статус", RED)

    add_text(
        slide,
        Inches(0.75), Inches(6.35), Inches(12.0), Inches(0.35),
        "Целевая аудитория: оператор смены, технический администратор, владелец или менеджер кафе.",
        12, False, MUTED
    )


def slide_alerts(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_bar(slide, "Пороги и алерты", "Правила формирования предупреждений и критических уведомлений")

    add_panel(slide, Inches(0.4), Inches(1.2), Inches(12.45), Inches(5.95), None)
    headers = ["Метрика", "Норма", "Warn", "Critical", "Кому", "Реакция"]
    widths = [Inches(2.0), Inches(1.15), Inches(1.2), Inches(1.3), Inches(2.1), Inches(4.0)]
    rows = [
        ["Доступность", ">=99.5%", "<99.5%", "<99.0%", "админ, владелец", "проверка API, БД, сети"],
        ["P95", "<=700 мс", "700-1500", ">1500", "техподдержка", "анализ медленных запросов"],
        ["Ошибки", "<1%", "1-3%", ">3%", "dev + support", "rollback или hotfix"],
        ["Печать", ">=98%", "95-98%", "<95%", "оператор, админ", "проверка принтера и очереди"],
        ["Задержки", "<10%", "10-20%", ">20%", "менеджер смены", "проверка кухни и маршрута"],
        ["Остатки", "0-2 SKU", "3-5", ">5", "кладовщик", "дозакуп и блокировка блюд"],
        ["MTTR", "<=15 мин", "15-30", ">30", "руководитель", "пересмотр регламента"],
    ]
    add_table_like(slide, Inches(0.6), Inches(1.5), widths, headers, rows)

    add_text(
        slide,
        Inches(0.75), Inches(6.38), Inches(12.0), Inches(0.28),
        "Алерт формируется при пересечении порога в двух подряд окнах измерения, чтобы снизить ложные срабатывания.",
        11, False, MUTED
    )


def slide_value(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_bar(slide, "Ценность метрик и итог", "Как мониторинг снижает риски и поддерживает SLA")

    add_panel(slide, Inches(0.55), Inches(1.35), Inches(6.0), Inches(4.7), "Какие риски снижаются")
    add_bullets(
        slide,
        Inches(0.85), Inches(1.85), Inches(5.45), Inches(3.8),
        [
            "остановка продаж из-за недоступности кассы",
            "потеря заказов из-за ошибок API и сбоев печати",
            "срыв обслуживания из-за задержек на кухне",
            "невозможность готовить блюда из-за дефицита ингредиентов",
        ],
        font_size=19,
    )

    add_panel(slide, Inches(6.8), Inches(1.35), Inches(5.95), Inches(4.7), "Как это влияет на SLA")
    add_bullets(
        slide,
        Inches(7.1), Inches(1.85), Inches(5.35), Inches(3.8),
        [
            "ускоряется обнаружение инцидентов",
            "снижается MTTR",
            "растёт предсказуемость работы смены",
            "повышается качество обслуживания клиента",
        ],
        font_size=19,
    )

    add_panel(slide, Inches(0.55), Inches(6.2), Inches(12.2), Inches(0.85), "Вывод")
    add_text(
        slide,
        Inches(0.8), Inches(6.48), Inches(11.6), Inches(0.28),
        "Мониторинг CafeHelp должен контролировать не только сервер, но и критичные бизнес-сценарии: заказ, печать, смену и склад. Это превращает мониторинг в управленческий инструмент, а не просто в набор технических графиков.",
        13, False, TEXT
    )


def slide_pitch(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RGBColor(239, 246, 255))
    add_top_bar(slide, "Краткий питч", "Финальный слайд для устной защиты")

    quote = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(1.55), Inches(11.7), Inches(4.7))
    quote.fill.solid()
    quote.fill.fore_color.rgb = WHITE
    quote.line.color.rgb = ACCENT
    quote.line.width = Pt(2)

    add_text(
        slide,
        Inches(1.2), Inches(1.95), Inches(11.0), Inches(3.9),
        "Я разработал систему мониторинга для сервиса CafeHelp, который автоматизирует работу кафе: приём заказов, смены, склад и печать чеков. В основу мониторинга я заложил не только технические метрики, такие как доступность, отклик, ошибки и MTTR, но и операционные показатели: успешность печати, долю задержанных заказов и критичные остатки на складе. Такой подход позволяет заранее обнаруживать риски, сокращать простой и поддерживать целевой SLA на уровне, достаточном для реальной работы кафе.",
        19, False, TEXT
    )

    add_text(slide, Inches(1.2), Inches(6.45), Inches(5.2), Inches(0.25), "CafeHelp Monitoring", 13, True, NAVY)


def main():
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(
        prs,
        "Система метрик и мониторинга для сервиса CafeHelp",
        "Задание 3. Измерение качества ИТ-сервиса, система контроля, дашборд и алерты",
    )
    slide_goal(prs)
    slide_metrics(prs)
    slide_method(prs)
    slide_dashboard(prs)
    slide_alerts(prs)
    slide_value(prs)
    slide_pitch(prs)

    prs.save(OUTPUT_PATH)
    print(OUTPUT_PATH)


if __name__ == "__main__":
    main()
